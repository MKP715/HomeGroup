"""
Update Draft_Foundation_Heritage_Humility.pptx with missing content from index.html.
- Enriches presenter notes with comprehensive talking points
- Adds missing body content to slides with significant gaps
- Adds slide transitions for smoother flow
- Adds section navigation hyperlinks
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

PPTX_PATH = r'c:/Users/MKP715/HomeGroup/HomeGroup/Draft_Foundation_Heritage_Humility.pptx'
HTML_PATH = r'c:/Users/MKP715/HomeGroup/HomeGroup/index.html'

# Colors matching the existing design
ACCENT = RGBColor(0xD4, 0xA8, 0x4B)  # Gold
BG_DARK = RGBColor(0x0D, 0x1B, 0x2A)
BG_CARD = RGBColor(0x13, 0x29, 0x3D)
TEXT_WHITE = RGBColor(0xE8, 0xE8, 0xE8)
TEXT_MUTED = RGBColor(0x99, 0xAA, 0xBB)
SUCCESS = RGBColor(0x4C, 0xAF, 0x50)
DANGER = RGBColor(0xE7, 0x4C, 0x3C)
WARNING = RGBColor(0xFF, 0x98, 0x00)

def extract_html_notes(html_path):
    """Extract data-notes from each slide in index.html."""
    with open(html_path, 'r', encoding='utf-8') as f:
        html = f.read()
    # Match data-notes="..." across the file
    notes = re.findall(r'data-notes="(.*?)">', html, re.DOTALL)
    # Decode HTML entities
    decoded = []
    for note in notes:
        note = note.replace('&mdash;', '—').replace('&amp;', '&')
        note = note.replace('&bull;', '•').replace('&rarr;', '→')
        note = note.replace('&ndash;', '–').replace('&frac12;', '½')
        note = note.replace('&ntilde;', 'ñ')
        note = re.sub(r'&#\d+;', '', note)  # Remove remaining numeric entities
        note = note.replace("'", "'").replace("'", "'")
        decoded.append(note.strip())
    return decoded


def add_text_box(slide, left, top, width, height, text, font_name='Calibri',
                 font_size=Pt(11), font_color=TEXT_WHITE, bold=False, italic=False,
                 alignment=PP_ALIGN.LEFT, word_wrap=True):
    """Helper to add a text box with formatted text to a slide."""
    from pptx.util import Emu as _Emu
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.italic = italic
    return txBox


def add_multiline_text_box(slide, left, top, width, height, lines,
                            font_name='Calibri', font_size=Pt(10),
                            font_color=TEXT_MUTED, line_spacing=1.15):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (text, is_bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = font_size
        run.font.color.rgb = color or font_color
        run.font.bold = is_bold
        p.space_after = Pt(2)

    return txBox


def add_transitions(prs):
    """Add smooth slide transitions to all slides."""
    for slide in prs.slides:
        # Add fade transition using XML
        transition = slide._element.makeelement(
            qn('p:transition'),
            {
                'spd': 'med',
                'advClick': '1',
            }
        )
        fade = transition.makeelement(qn('p:fade'), {'thruBlk': '1'})
        transition.append(fade)
        # Insert transition before the first element or at the end
        slide._element.insert(0, transition)


def update_notes(prs, html_notes):
    """Update presenter notes, merging HTML notes where they're more comprehensive."""
    for i, slide in enumerate(prs.slides):
        if i >= len(html_notes):
            break

        html_note = html_notes[i]
        if not html_note:
            continue

        # Get existing notes
        existing_notes = ''
        if slide.has_notes_slide:
            existing_notes = slide.notes_slide.notes_text_frame.text

        # If HTML note has significantly more content, append it
        if len(html_note) > len(existing_notes) + 200:
            # HTML has substantially more - replace with merged version
            merged = existing_notes.strip()
            merged += "\n\n--- ADDITIONAL TALKING POINTS (from interactive version) ---\n\n"
            merged += html_note
            if not slide.has_notes_slide:
                slide.notes_slide  # Creates notes slide
            # Clear and rewrite
            notes_tf = slide.notes_slide.notes_text_frame
            notes_tf.clear()
            p = notes_tf.paragraphs[0]
            p.text = merged
            p.font.size = Pt(10)
        elif len(html_note) > 500 and not slide.has_notes_slide:
            # No existing notes, HTML has content
            notes_tf = slide.notes_slide.notes_text_frame
            notes_tf.clear()
            p = notes_tf.paragraphs[0]
            p.text = html_note
            p.font.size = Pt(10)


def add_missing_content_slide7(slide):
    """Slide 7 (Service Positions) - Add checklist items and positions table."""
    # Add service checklist as compact text
    checklist_text = (
        "What AA Group Members Do (P-16, pp. 18-19):\n"
        "• Provide and maintain a meeting place\n"
        "• Arrange programs for meetings\n"
        "• Collect and allocate Seventh Tradition contributions\n"
        "• Maintain Conference-approved literature\n"
        "• Provide Grapevine/La Viña materials\n"
        "• Offer refreshments\n"
        "• Assist alcoholics in finding meetings\n"
        "• Answer calls for help\n"
        "• Air and resolve group problems\n"
        "• Sustain contact with A.A. locally and internationally"
    )
    add_text_box(slide,
                 Inches(0.5), Inches(2.8), Inches(4.5), Inches(2.5),
                 checklist_text, font_size=Pt(9), font_color=TEXT_MUTED)

    # Add positions table as text
    positions_text = (
        "Core Service Positions (P-16, pp. 19-27):\n"
        "Chairperson — Coordinate activities (1+ yr sobriety)\n"
        "Secretary — Records, announcements (6 mo–1 yr)\n"
        "Treasurer — Manage funds/records (1–2+ yrs)\n"
        "GSR — Link to General Service Conference (2–3 yrs)\n"
        "Intergroup Rep — Link to central office (1–2 yrs)\n\n"
        "Additional: Grapevine Rep, Literature Rep, CPC, PI,\n"
        "Corrections, Treatment, Accessibilities liaisons"
    )
    add_text_box(slide,
                 Inches(5.2), Inches(2.8), Inches(4.5), Inches(2.5),
                 positions_text, font_size=Pt(9), font_color=TEXT_MUTED)

    # Add Living Sober quote
    quote_text = '"In A.A., no one is \'above\' or \'below\' anyone else." — Living Sober, p. 15'
    add_text_box(slide,
                 Inches(0.5), Inches(5.0), Inches(9.0), Inches(0.4),
                 quote_text, font_size=Pt(9), font_color=ACCENT, italic=True)


def add_missing_content_slide13(slide):
    """Slide 13 (Key Quotations) - Add missing quotes."""
    missing_quotes = [
        ('"We have a way out on which we can absolutely agree, and upon which we can join in '
         'brotherly and harmonious action."\n— Big Book, p. 17'),
        ('"No satisfaction has been deeper and no joy greater than in a Twelfth Step job well done. '
         'To watch the eyes of men and women open with wonder as they move from darkness into '
         'light..."\n— As Bill Sees It, p. 29'),
        ('"Abandon yourself to God as you understand God. Admit your faults to Him and to your '
         'fellows. Clear away the wreckage of your past. Give freely of what you find and join us."\n'
         '— Big Book, p. 164 — The Great Invitation'),
        ('"It is the great paradox of A.A. that we know we can seldom keep the precious gift of '
         'sobriety unless we give it away."\n— 12&12, p. 151'),
    ]
    y_pos = Inches(3.6)
    for quote in missing_quotes:
        add_text_box(slide,
                     Inches(0.5), y_pos, Inches(9.0), Inches(0.55),
                     quote, font_size=Pt(8.5), font_color=TEXT_MUTED, italic=True)
        y_pos += Inches(0.55)


def add_missing_content_slide14(slide):
    """Slide 14 (Literature Guide & Closing) - Add lit reference + closing."""
    lit_ref = (
        "COMPLETE LITERATURE REFERENCE:\n"
        "• Big Book: pp. 17, 89, 97, 100, 164 — Fellowship, service, invitation\n"
        "• 12 Steps & 12 Traditions: pp. 106-125, 129-154, 189 — Step 12, Traditions\n"
        "• P-16 'The A.A. Group': pp. 12-16, 18-30 — Definition, membership, service\n"
        "• P-15 'Sponsorship Q&A': pp. 7-8, 12, 23-26 — Home group connection\n"
        "• As Bill Sees It: pp. 9, 29, 50, 90, 117 — Groups, service, loneliness\n"
        "• AA Service Manual: Ch. 1-2, pp. S1, S15, S25-31 — Structure, GSR role\n"
        "• Living Sober: pp. 12-17, 75-81 — Getting active, trying meetings\n"
        "• AA Comes of Age: pp. 99, 101, 163 — Group conscience, history\n"
        "• Dr. Bob & Good Oldtimers: pp. 101, 144, 167, 338 — Early groups\n"
        "• Grapevine 'Home Group: Heartbeat' — Stories, service, traditions\n"
        "• Twelve Concepts: Concept I — Ultimate authority in Fellowship"
    )
    add_text_box(slide,
                 Inches(0.3), Inches(1.8), Inches(9.4), Inches(2.5),
                 lit_ref, font_size=Pt(8.5), font_color=TEXT_MUTED)

    # Add closing activity
    closing = (
        "CLOSING ACTIVITY — Personal Commitment:\n"
        "What is ONE action you will take THIS WEEK to strengthen your home group connection?\n"
        "If you don't have a home group, what meeting will you visit this week?\n\n"
        '"I am responsible. When anyone, anywhere, reaches out for help,\n'
        'I want the hand of A.A. always to be there. And for that: I am responsible."\n'
        '— The Responsibility Statement (1965)'
    )
    add_text_box(slide,
                 Inches(0.3), Inches(4.2), Inches(9.4), Inches(1.3),
                 closing, font_size=Pt(9), font_color=ACCENT, italic=True)


def add_missing_content_slide15(slide):
    """Slide 15 (Bonus Scenarios) - Add missing scenarios."""
    additional = (
        "ADDITIONAL SCENARIOS:\n\n"
        "The Money Question: Group has $3,000 in treasury. Members disagree on whether to donate "
        "to intergroup, keep a larger prudent reserve, or split between district, NETA 65, and GSO. "
        "What does Tradition 7 say? What is a 'prudent reserve'?\n\n"
        "The Online Question: Since the pandemic, group has Zoom option. Should it be permanent? "
        "Some say 'real AA' is in person only. What does Tradition 4 say about autonomy?\n\n"
        "The Anonymity Break: A member posts a meeting photo on social media, tagging others by "
        "name. Some upset, others say 'relax.' What do Traditions 11 and 12 say?\n\n"
        "The Crosstalk Debate: Group has 'no crosstalk' guideline. A newer member says it feels "
        "cold — wants real dialogue. Old-timer says it protects vulnerable members. Is it a Tradition "
        "or a guideline? How does Tradition 4 apply?\n\n"
        "The Dying Group: Once had 30 members, now averages 6. All old-timers, no newcomers. "
        "Lease is expensive. Should they close? What about the alcoholic who might walk in next week?"
    )
    add_text_box(slide,
                 Inches(0.3), Inches(2.5), Inches(9.4), Inches(3.0),
                 additional, font_size=Pt(8), font_color=TEXT_MUTED)


def add_missing_content_slide23(slide):
    """Slide 23 (Fun Facts) - Add missing trivia."""
    additional = (
        "MORE TRIVIA:\n\n"
        "What is 'Jake Leg' and who at Dallas Central Office had it?\n"
        "→ Jamaica Ginger Paralysis — caused by toxic Prohibition-era alcohol. Dick P., "
        "the first Central Office manager, suffered from it.\n\n"
        "What article quadrupled AA membership — and directly led to Dallas AA's founding?\n"
        "→ Jack Alexander's Saturday Evening Post article (March 1, 1941). Esther E.'s "
        "husband read it, gave her an ultimatum, and she contacted GSO."
    )
    add_text_box(slide,
                 Inches(0.3), Inches(3.5), Inches(9.4), Inches(2.0),
                 additional, font_size=Pt(9), font_color=TEXT_MUTED)


def add_missing_content_slide26(slide):
    """Slide 26 (Why Anonymity) - Add 'What Anonymity Is NOT' content."""
    what_not = (
        "What Anonymity Is NOT:\n\n"
        "✗ NOT secrecy within the fellowship — we share openly within AA. "
        "First names are for humility, not hiding.\n\n"
        "✗ NOT that you can't tell anyone — many members are open with family, friends, "
        "employers. Anonymity is about PUBLIC media, not private conversations.\n\n"
        "✗ NOT a loophole for bad behavior — group conscience can address misconduct. "
        "Anonymity protects AA's public reputation, not individual behavior within the fellowship.\n\n"
        '"The spiritual substance of anonymity is sacrifice." — Bill W., AA Grapevine, Jan 1955'
    )
    add_text_box(slide,
                 Inches(0.3), Inches(3.2), Inches(9.4), Inches(2.3),
                 what_not, font_size=Pt(9), font_color=TEXT_MUTED)


def add_missing_content_slide36(slide):
    """Slide 36 (Anonymity Closing) - Add literature table and action items."""
    lit_table = (
        "ESSENTIAL LITERATURE ON ANONYMITY:\n"
        "• 12&12: Traditions 11 (pp. 180-183) & 12 (pp. 184-187)\n"
        "• Big Book: Foreword, 'Working With Others' (pp. 89-103)\n"
        "• 'Understanding Anonymity' Pamphlet (22 pages)\n"
        "• AA Comes of Age: Anonymity chapter (pp. 135-141)\n"
        "• 'Anonymity Online' GSO Bulletin (available at aa.org)\n"
        "• The Language of the Heart: pp. 67-72\n"
        "• Pass It On: pp. 282-283, 335-337"
    )
    add_text_box(slide,
                 Inches(0.3), Inches(3.0), Inches(4.5), Inches(2.0),
                 lit_table, font_size=Pt(8.5), font_color=TEXT_MUTED)

    actions = (
        "ACTION ITEMS — Bring It Back:\n\n"
        "This Week:\n"
        "• Review your social media for potential breaks\n"
        "• Read Traditions 11 & 12\n"
        "• Discuss anonymity with your sponsor\n\n"
        "This Month:\n"
        "• Propose Traditions meeting on anonymity\n"
        "• Suggest technology policy if your group lacks one\n"
        "• Volunteer for PI committee\n\n"
        "This Year:\n"
        "• Serve without seeking credit\n"
        "• Place principles before personalities\n"
        "• Model anonymity for those you sponsor"
    )
    add_text_box(slide,
                 Inches(5.0), Inches(3.0), Inches(4.7), Inches(2.5),
                 actions, font_size=Pt(8.5), font_color=TEXT_MUTED)


def main():
    print("Loading PPTX...")
    prs = Presentation(PPTX_PATH)

    print("Extracting HTML notes...")
    html_notes = extract_html_notes(HTML_PATH)
    print(f"  Found {len(html_notes)} HTML slide notes")

    print("Updating presenter notes...")
    update_notes(prs, html_notes)

    print("Adding missing content to slides...")

    slides = list(prs.slides)

    # Slide 7 (index 6) - Service Positions: add checklist + positions table
    print("  Slide 7: Adding service checklist & positions table")
    add_missing_content_slide7(slides[6])

    # Slide 13 (index 12) - Key Quotations: add missing quotes
    print("  Slide 13: Adding missing quotes")
    add_missing_content_slide13(slides[12])

    # Slide 14 (index 13) - Literature Guide & Closing
    print("  Slide 14: Adding literature reference & closing")
    add_missing_content_slide14(slides[13])

    # Slide 15 (index 14) - Bonus Scenarios: add missing scenarios
    print("  Slide 15: Adding missing scenarios")
    add_missing_content_slide15(slides[14])

    # Slide 23 (index 22) - Fun Facts: add missing trivia
    print("  Slide 23: Adding missing trivia questions")
    add_missing_content_slide23(slides[22])

    # Slide 26 (index 25) - Why Anonymity: add "What Anonymity Is NOT"
    print("  Slide 26: Adding 'What Anonymity Is NOT' content")
    add_missing_content_slide26(slides[25])

    # Slide 36 (index 35) - Anonymity Closing: add lit table + action items
    print("  Slide 36: Adding literature & action items")
    add_missing_content_slide36(slides[35])

    print("Adding slide transitions...")
    add_transitions(prs)

    print("Saving updated PPTX...")
    prs.save(PPTX_PATH)
    print(f"Done! Saved to {PPTX_PATH}")


if __name__ == '__main__':
    main()
